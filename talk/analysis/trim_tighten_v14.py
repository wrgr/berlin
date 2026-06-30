"""Trim + tighten Calibrate_the_Humans_v14.pptx -> 21 main slides + 3 backup.
- Merge S14+S15: retitle S14 to cover exploration+grammar, fold grammar into its 3 columns
  (keeps S14's figure); S15 (grammar detail) -> backup.
- S16 (RF/PCA/motif) -> backup.
- Merge S21+S22: retitle/condense S21, fold S22's students-as-method point; S22 -> backup.
- Tighten ~17 of the wordiest bullets across the deck.
Backup slides are relocated after 'Thank You' (no deletion => no orphaned media => no bloat).
Operates in place on talk/Calibrate_the_Humans_v14.pptx.
"""
from pathlib import Path
from pptx import Presentation
TALK=Path(__file__).resolve().parent.parent
DECK=TALK/"Calibrate_the_Humans_v14.pptx"
p=Presentation(str(DECK))
S=lambda n:p.slides[n-1]

def settext(shape,text):
    tf=shape.text_frame; p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=text
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        p0.add_run().text=text
    for para in tf.paragraphs[1:]: para._p.getparent().remove(para._p)

def find(slide,snip):
    for sh in slide.shapes:
        if sh.has_text_frame and snip in sh.text_frame.text: return sh
    return None

def setby(n,snip,newtext):
    sh=find(S(n),snip)
    if sh is None: print("  !! no match S%d %r"%(n,snip[:30])); return
    settext(sh,newtext)

# ---------- merge S14 + S15 ----------
setby(14,"Behavioral Mechanism","The Language of Proofreading: Exploration & Grammar")
setby(14,"Camera Rotation","Exploration")
setby(14,"Experts accumulate ~2.18","~2.18× more camera rotation (1014° vs 466°), more viewpoints, faster tempo.")
setby(14,"Viewpoints & Events","Action Mix")
setby(14,"Experts register more rotations","Experts vs proto-experts split navigate / segment / annotate differently.")
setby(14,"Tempo","Grammar")
setby(14,"Experts move faster","Navigate↔segment transitions differ by cohort — experts edit in sustained runs (figure in backup).")
setby(15,"The Language of Proofreading: Action","Backup — Action Mix & Grammar (detail)")
setby(16,"What Separates Experts","Backup — What Separates Experts (RF / PCA / motif)")

# ---------- merge S21 + S22 ----------
setby(21,"One Node in an Interconnected","Where This Lives — Students as the Method")
setby(21,"High-throughput Integrative Mouse","HI-MC connectome (Lichtman UM1) — the map this work serves.")
setby(21,"Learning Engineering @ JHU School","Learning Engineering @ JHU School of Education — trains the next workforce.")
setby(21,"Automated segmentation, proofreading,","Automated segmentation, proofreading, and inference — one integrated agenda.")
setby(21,"Calibrating students to succeed in","MICrONS relied on this workforce; students became co-authors (Daniel Xenes now leads NeuVue). Calibrated student judgment is a replicable method.")
setby(22,"Outreach: Students as a Method","Backup — Outreach (detail)")

# ---------- tighten exposition ----------
TIGHTEN=[
 (2,"Engineering is problem-solving under real-world limits","Engineering is problem-solving under real-world limits — many threads converge on one challenge."),
 (5,"Connectivity proxies and completeness metrics","Completeness metrics (i2g, NRI, ERL, ARI) — measured against the dataset, edge-agnostic (incl. our own)."),
 (5,"Pathfinder, Focused Proofreading, Autoproof","Pathfinder, Focused Proofreading, NEURD, RoboEM — reach 'enough' faster, but don't decide when."),
 (5,"RoboEM (Motta et al.) validated automation","RoboEM and our motif work — validated by analysis outcome. Task-relative, but not yet operational."),
 (6,"Humans are messy","Humans are messy — they tire, vary in focus, and bring judgment that's hard to formalize."),
 (6,"also have biases","Machine agents have their own biases and heavy resource needs — no single agent wins everywhere."),
 (6,"The optimization framework is largely agnostic","The framework is agnostic to human vs machine — what matters is matching capability to decision."),
 (7,"Fix merges, splits, and missed partners","Fix merges, splits, missed partners — active editing to improve accuracy."),
 (7,"CONFIRMS /","CONFIRMS / MICrONS-style — certify, don't edit; confirm quality is sufficient for the claim."),
 (8,"To allocate judgment, you have to price it","To allocate judgment you must price it — per-task and per-region competence, and agreement."),
 (11,"Promotion judged by hand, after the fact","Promotion judged by hand, after the fact — expensive expert time spent watching."),
 (11,"Behavior and the grammar of proofreading give formative","Behavior + grammar give formative signals in training, summative ones at promotion."),
 (12,"Calibrated annotators converge to expert-level agreement","Calibrated annotators converge to expert-level agreement — calibration works."),
 (12,"Behavioral signals cleanly separate experts","Behavior separates experts from non-experts — skill leaves a measurable trace."),
 (19,"36 novice proofreaders","36 novice proofreaders — JHU undergraduates (26 founders; 3 weeks training, then part/full-time)."),
 (20,"Students whose decisions agreed with experts were promoted","Students who agreed with experts were promoted to a proto-expert tier with expert-task write access."),
 (23,"The missing piece is not better algorithms","The missing piece isn't better algorithms — it's a principled price for human judgment."),
]
for n,snip,new in TIGHTEN: setby(n,snip,new)

# ---------- relocate S15, S16, S22 to backup (after Thank You) ----------
lst=p.slides._sldIdLst; ids=list(lst)
move=[ids[14],ids[15],ids[21]]   # 0-indexed: S15,S16,S22
for el in move: lst.remove(el)
for el in move: lst.append(el)

p.save(str(DECK))
print("saved %s — %d slides (21 main + 3 backup)"%(DECK.name,len(p.slides._sldIdLst)))
