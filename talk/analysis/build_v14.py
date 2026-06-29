"""build_v14 — apply the honest-reframe to the rebaselined 'Calibrate the Humans' deck.
Base: the user's rebaselined CalibratetheHumansNotJusttheMicroscope.pptx (24 slides).
Changes (lead with the robust, model-free signal; label the engineered AUCs an exploratory ceiling;
correct the rho=-0.44 framing to a selection artifact, NOT 'style != proficiency'):
  S13  swap tier-AUC chart -> fig_expertise_evidence.png; reword headline bullet
  S14  '~2x' -> '~2.18x (1014 vs 466 deg)', model-free
  S16  subtitle: mark EXPLORATORY (post-hoc features)
  S17  add a selection-artifact + data-limitation caveat (calibrated cohort, no true novices)
Outputs talk/berlin_deck_v14.pptx.
"""
import os, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.text.text import _Paragraph

SCR=Path(os.environ.get("BERLIN_SCRATCH","/tmp/claude-0/-home-user-berlin/eb94da13-135d-5ec6-9497-7a16104e77d6/scratchpad"))
TALK=Path(os.environ.get("BERLIN_TALK","/home/user/berlin/talk"))
BASE=SCR/"deck_rebaseline.pptx"
OUT=TALK/"berlin_deck_v14.pptx"
p=Presentation(str(BASE))

def shape_containing(slide, sub):
    for sh in slide.shapes:
        if sh.has_text_frame and sub in sh.text_frame.text:
            return sh
    raise KeyError(sub)

def set_para_text(para, text):
    if not para.runs:
        para.add_run().text=text; return
    para.runs[0].text=text
    for r in para.runs[1:]:
        r._r.getparent().remove(r._r)

def replace_single_block(slide, sub, newtext):
    sh=shape_containing(slide, sub); tf=sh.text_frame
    set_para_text(tf.paragraphs[0], newtext)
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)

def append_caveat(slide, sub, newtext):
    """clone the last paragraph (keeps bullet/font), set caveat text."""
    sh=shape_containing(slide, sub); tf=sh.text_frame
    last=tf.paragraphs[-1]._p
    new=copy.deepcopy(last); last.getparent().append(new)
    set_para_text(_Paragraph(new, tf.paragraphs[-1]._parent), newtext)

def swap_picture(slide, width_in, path, tol=0.3):
    for sh in list(slide.shapes):
        if sh.shape_type==13 and abs(Emu(sh.width).inches-width_in)<tol:
            l,t,w,h=sh.left,sh.top,sh.width,sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(str(path),l,t,w,h)
            return True
    raise KeyError("no picture ~%.1fin"%width_in)

S=lambda n:p.slides[n-1]

# --- S13: honest evidence headline + chart swap ---
replace_single_block(S(13), "Behavior separates experts from proto-experts across feature sets",
    "Experts explore ~2.18× more — the rawest signal, no model. Across tiers (cross-validated), naive "
    "4-count features and the learned 10-motif dictionary both separate experts from proto-experts at "
    "AUC 0.81; the 28-feature designed bank reaches 0.98 but is engineered post-hoc on n=16 — an "
    "exploratory ceiling. CV rules out trivial fit: 28 pure-noise features collapse to 0.45.")
swap_picture(S(13), 5.3, TALK/"fig_expertise_evidence.png")

# --- S14: precise, model-free rotation ---
replace_single_block(S(14), "Experts accumulate ~2× more total camera rotation",
    "Experts accumulate ~2.18× more total camera rotation (1014° vs 466°), surveying the structure "
    "from significantly more viewpoints before making decisions — a simple, model-free difference.")

# --- S16: mark exploratory ---
replace_single_block(S(16), "RANDOM FOREST IMPORTANCE",
    "RANDOM FOREST IMPORTANCE · DESIGNED FEATURES · N=16 · EXPLORATORY (POST-HOC FEATURES)")

# --- S17: selection-artifact + data-limitation caveat ---
append_caveat(S(17), "Promoted ≈ Expert < Unpromoted",
    "Caveat (n=16): this accuracy cohort is calibrated — experts + agreement-promoted, no true "
    "novices — so style tags the cohort but cannot rank accuracy. The ρ=−0.44 'expertise-vs-accuracy' "
    "dip is a selection artifact (promotion was gated on grader-agreement, the very metric), not "
    "'style ≠ proficiency.'")

p.save(str(OUT))
print("saved", OUT, "slides:", len(p.slides))
