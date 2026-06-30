#!/usr/bin/env python3
"""Complete the speaker notes on berlin_deck_v11.pptx.

The deck already carries rich first-person speaker notes on the argument slides (1-13, 19-21);
this fills the evidence slides (14-18) and the close/backup (22-23), which shipped as one-liners
or empty. Notes are spoken prose in the voice of the existing slides 5/7 notes — landing points,
not a teleprompter. Source of truth for the beats is ../talk_script.md; the numbers match
../figure_descriptions.md.

Idempotent and reversible: re-running just re-sets the same notes (git tracks the binary).
Only the slides in NOTES are touched; every other slide's notes are left untouched.
Run with BERLIN_TALK pointing at talk/ (defaults to the repo-relative talk/ dir).
"""
import os
from pathlib import Path
from pptx import Presentation

HERE = Path(__file__).resolve().parent
TALK = Path(os.environ.get("BERLIN_TALK", HERE.parent))
DECK = TALK / "berlin_deck_v11.pptx"

# 1-based slide number -> speaker notes
NOTES = {
    14: (
        "This is the mechanism behind that AUC — what experts actually do differently. Four "
        "measures on the dense telemetry, eight per group. Experts accumulate about twice the total "
        "camera rotation — roughly a thousand degrees versus four-fifty — and inspect each neuron "
        "from more than twice as many distinct viewpoints, about 14 versus 6. They also do more "
        "overall and move a bit faster. So expertise isn't a secret edit you make or don't — it's a "
        "3-D exploration style: experts look harder, from more angles, before they commit. The "
        "metric is per-differstack rotation, so this isn't an artifact of cross-session camera "
        "jumps. Landing line: skill is in the *how*, and the *how* is visual inspection."
    ),
    15: (
        "Beyond how much they explore — what they do, and in what order. Two things. The action "
        "mix: experts spend more of their events in segmentation; proto-experts navigate more. And "
        "the grammar — the transition probabilities between navigate and segment. Experts show "
        "stronger segment-to-segment persistence, about 0.69 versus 0.60: once they start editing "
        "they stay in it — sustained runs. Proto-experts hop: navigate, edit, navigate. This is the "
        "JIGSAWS 'language of surgery' analogy made literal — there's a grammar to proofreading, and "
        "fluency in it is what expertise looks like. (If asked: this panel is descriptive cohort "
        "means; 'annotate' is near zero here because annotation lives in other task types.)"
    ),
    16: (
        "Three independent views of the same question — what separates the cohorts — so you can see "
        "the signal isn't an artifact of one method. Left: RandomForest importance over the designed "
        "features — the top features are the 3-D exploration kinematics, rotation and viewpoints. "
        "Center: a 2-D PCA of the standardized features — cohorts separate along the first "
        "component, and the promoted students (circled) sit right among the experts. Right: usage of "
        "the ten unsupervised learned motifs — proto-experts over-use the navigate-dominated motifs, "
        "experts spread into segment-and-rotation ones. A behavioral dialect. Landing: hand-built, "
        "dimensionality-reduced, and fully unsupervised all point to the same place — it's real."
    ),
    17: (
        "This is the pivot of the whole evidence arc — a negative and a positive together. Left "
        "panel, distance-to-ground-truth for three groups: experts 309, promoted students 312, "
        "unpromoted 460. Promotion worked — it selected people who perform at expert level — which "
        "means achieved accuracy is ceiling-compressed; you literally can't use it to rank the "
        "calibrated workforce. That's the negative. The positive is the right panel: even after that "
        "convergence, which individual *decisions* are risky stays readable. Bin tasks by how slow "
        "they were for that person — the slowest quintile fails about 2.2x more than the fastest, "
        "with no ground truth. Landing line, and it's the thesis: don't rank people — flag risky "
        "decisions."
    ),
    18: (
        "The capstone — the risk axis from slide 7 made operational. From a task's annotation-"
        "category structure alone, no ground truth, we predict whether it'll contain an error at "
        "AUC 0.76 on held-out cells — honest grouped-by-cell cross-validation, against a permutation "
        "null of 0.47, p below 0.001. I want to be straight about one thing: an earlier 0.92 was "
        "cell-identity leakage — only 28 benchmark cells, so random CV lets the model memorize the "
        "cell. Grouping fixes it. What survives is largely 'flag the few intrinsically hard cells' — "
        "three killer cells above 78% error — and per-person competence within a fixed cell stays "
        "weak, about 0.55. That's consistent: per-decision, not per-person. Landing: you can score "
        "risk *before* you spend a minute of expert time — exactly what impact-times-risk needs."
    ),
    22: (
        "Let me land it. Proofreading a connectome isn't a quality problem you finish — it's an "
        "optimization problem under a hard constraint: limited, noisy, expensive expert judgment. "
        "The piece that's been missing is the *price* of that judgment — and once you can price it, "
        "per-decision and ground-truth-free, you can allocate it where impact and risk are highest. "
        "And the same measurement that prices the budget is the one that develops the workforce. So "
        "the line I want you to leave with: calibrate the people, not just the microscope. Thank you."
    ),
    23: (
        "Only if someone pushes on 'can't you just rank your proofreaders.' We pressure-tested that "
        "hard, and the answer is no — and the careful version matters. Per-annotator accuracy is not "
        "predictable: not on the ceiling-bound label accuracy, not on the variance-rich distance-to-"
        "ground-truth; no scale transform and no model — Ridge, RandomForest, gradient boosting, "
        "kNN — recovers it, best Spearman 0.26 at p=0.25. And the 0.14 leave-one-out AUC isn't "
        "'worse than chance,' it's inside the null, 0.45 plus-or-minus 0.20. What does survive is "
        "per-decision flagging, AUC 0.59, robust to task size — 3-D structural difficulty is the "
        "follow-up. So the honest line is the same as the main line: flag decisions, don't rank "
        "people."
    ),
}


def main():
    if not DECK.exists():
        raise SystemExit(f"deck not found: {DECK}")
    prs = Presentation(str(DECK))
    n = len(prs.slides)
    for num, text in sorted(NOTES.items()):
        if not (1 <= num <= n):
            print(f"!! slide {num} out of range (deck has {n}); skipping")
            continue
        slide = prs.slides[num - 1]
        before = len(slide.notes_slide.notes_text_frame.text.strip()) if slide.has_notes_slide else 0
        slide.notes_slide.notes_text_frame.text = text
        print(f"slide {num:2d}: notes {before:4d} -> {len(text):4d} chars")
    prs.save(str(DECK))
    print(f"saved {DECK}")


if __name__ == "__main__":
    main()
