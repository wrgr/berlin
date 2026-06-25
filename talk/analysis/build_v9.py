"""Deck v8 -> v9: thread 'per-decision, not per-person' through the main line (slides 12 & 17) and
tighten the size-vs-difficulty wording on the backup slide (we controlled task SIZE, not 3-D
difficulty). Reversible/scripted. Run with BERLIN_TALK pointing at talk/."""
import os
from pathlib import Path
from pptx import Presentation
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v8.pptx")); S=prs.slides
def shape_in(sl,sub):
    for sh in sl.shapes:
        if sh.has_text_frame and sub in sh.text_frame.text: return sh
def set_text(sh,new):
    p=sh.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text=new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else: p.add_run().text=new
    for ex in sh.text_frame.paragraphs[1:]: ex._p.getparent().remove(ex._p)
log=[]
# Slide 17 — make the result the headline
set_text(shape_in(S[16],"Calibration converges"),"Competence is legible per-decision, not per-person"); log.append("S17 retitle")
# Slide 12 — thread the phrase into the closing thesis
set_text(shape_in(S[11],"Convergence proves the calibration works"),
    "Convergence proves the calibration works; separation proves expertise is predictable — but competence reads per-decision, not per-person: flag risky decisions, don't rank people."); log.append("S12 thesis thread")
# Slide 22 (backup) — task SIZE not difficulty; add transform/model robustness
set_text(shape_in(S[21],"Per-ANNOTATOR accuracy"),
    "Per-ANNOTATOR accuracy isn't predictable — not on ceiling label accuracy, not on variance-rich distance-to-GT; no scale transform or model (Ridge/RF/GBM/kNN) recovers it (best ρ=0.26, p=0.25), and the 0.14 is within the leave-one-out null (0.45±0.20). Per-DECISION flagging is the real signal (AUC 0.59; robust to task size — 3-D difficulty is a follow-up)."); log.append("S22 caption: size-not-difficulty + transform robustness")
prs.save(str(TALK/"berlin_deck_v9.pptx"))
print("saved berlin_deck_v9.pptx |",len(prs.slides),"slides | edits:",log)
