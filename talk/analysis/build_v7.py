"""Deck v6 -> v7: apply the three flagged recommendations.
1. Restore a closing slide (bookend of the title slide) so the thesis lands.
2. Slide 18: fix the 'Science & Engineering' label separator (1 space -> 3, matching other rows).
3. Slide 18: 'intersectional' -> 'interconnected'.
Scripted/reversible. Run with BERLIN_TALK pointing at talk/."""
import os, copy
from pathlib import Path
from pptx import Presentation
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v6.pptx")); S=prs.slides

def shape_in(slide, substr):
    for sh in slide.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text: return sh
    return None
def set_shape_text(shape, new):             # collapse to one line, keep first run's font
    tf=shape.text_frame; p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=new
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else: p0.add_run().text=new
    for extra in tf.paragraphs[1:]: extra._p.getparent().remove(extra._p)
def sub_run(shape, old, new):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old in r.text: r.text=r.text.replace(old,new); return True
    return False
def duplicate_slide(prs, index):            # clone a slide (text-only slide -> no rel fixups needed)
    src=prs.slides[index]
    new=prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes): sh._element.getparent().remove(sh._element)
    for sh in src.shapes: new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new

log=[]
# 1) closing slide — bookend of slide 1 (same layout/fonts), thesis restated
close=duplicate_slide(prs,0)
set_shape_text(shape_in(close,"Calibrate the students"), "Proofreading is an optimization problem.")
set_shape_text(shape_in(close,"When have we proofread enough"), "The missing piece is the price of human judgment.")
set_shape_text(shape_in(close,"Will Gray Roncal"), "Calibrate the people, not just the microscope.   ·   Thank you.")
log.append(f"added closing slide (now {len(prs.slides)} slides)")
# 2) slide 18 label separator: 1 space -> 3 (match other rows)
s18=prs.slides[17]
for sh in s18.shapes:
    if sh.has_text_frame and "Science" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text=="& Engineering ": r.text="& Engineering   "; log.append("S18 separator fixed")
# 3) slide 18 title wording
if sub_run(shape_in(s18,"intersectional ecosystem"),"intersectional","interconnected"): log.append("S18 intersectional->interconnected")
prs.save(str(TALK/"berlin_deck_v7.pptx"))
print("saved berlin_deck_v7.pptx |",len(prs.slides),"slides | edits:",log)
