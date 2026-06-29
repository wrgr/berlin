"""Merge S3 (Proofreaders as Agents) into S4 (Route the Work) -> one 'agents + routing' slide.
S4 (the impact x risk 2x2 matrix) is the base; fold S3's agent taxonomy into S4's intro line;
move S3 (with its agent diagram) to backup. 21 main -> 20 main, +1 backup. In place on the v14 deck.
"""
from pathlib import Path
from pptx import Presentation
DECK=Path(__file__).resolve().parent.parent/"Calibrate_the_Humans_v14.pptx"
p=Presentation(str(DECK)); S=lambda n:p.slides[n-1]

def settext(shape,text):
    tf=shape.text_frame; p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=text
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else: p0.add_run().text=text
    for para in tf.paragraphs[1:]: para._p.getparent().remove(para._p)
def setby(n,snip,new):
    for sh in S(n).shapes:
        if sh.has_text_frame and snip in sh.text_frame.text: settext(sh,new); return
    print("  !! no match S%d %r"%(n,snip[:30]))

# S4 becomes the merged slide
setby(4,"Route the Work to the Right Agent","Proofreaders as Agents — Route Each Decision to the Right One")
setby(4,"Experts, students, and ML are all agents in the loop",
      "From naive students catching obvious errors to experts handling the hard calls, experts, "
      "students, and ML are all agents in the loop — each with distinct capability, cost, and scope. "
      "The question is which agent each decision requires; routing well stretches a finite budget.")
# S3 (agent diagram) -> backup
setby(3,"Proofreaders as Agents","Backup — Proofreaders as Agents (detail)")
lst=p.slides._sldIdLst; el=list(lst)[2]; lst.remove(el); lst.append(el)

p.save(str(DECK))
print("saved %s — %d slides (20 main + 4 backup)"%(DECK.name,len(p.slides._sldIdLst)))
