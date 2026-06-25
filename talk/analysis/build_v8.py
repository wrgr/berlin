"""Deck v7 -> v8: add a BACKUP slide for the 'competence is per-decision, not per-person' result
(fig_accuracy_unpredictable.png). Clones the kinematics figure-slide for consistent styling, swaps
the image and text. Reversible/scripted. Run with BERLIN_TALK pointing at talk/."""
import os, copy
from pathlib import Path
from pptx import Presentation
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v7.pptx"))
def shape_in(sl,sub):
    for sh in sl.shapes:
        if sh.has_text_frame and sub in sh.text_frame.text: return sh
def set_text(sh,new):
    tf=sh.text_frame; p=tf.paragraphs[0]
    if p.runs:
        p.runs[0].text=new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else: p.add_run().text=new
    for ex in tf.paragraphs[1:]: ex._p.getparent().remove(ex._p)
def dup(prs,i):                                  # clone slide i, append at end
    src=prs.slides[i]; new=prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes): sh._element.getparent().remove(sh._element)
    for sh in src.shapes: new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new
new=dup(prs,13)                                  # slide 14 = kinematics (single wide figure + caption)
set_text(shape_in(new,"Behavioral mechanism"),"Backup — competence is legible per-decision, not per-person")
for sh in list(new.shapes):                      # swap the figure (drop the copied pic, add the new one)
    if sh.shape_type==13:
        l,t,w=sh.left,sh.top,sh.width; sh._element.getparent().remove(sh._element)
        new.shapes.add_picture(str(TALK/"fig_accuracy_unpredictable.png"),l,t,width=w); break
set_text(shape_in(new,"Experts accumulate"),
  "Per-ANNOTATOR accuracy isn't behaviorally predictable — on ceiling label accuracy OR variance-rich "
  "distance-to-GT (Ridge ρ=0.07). The 0.14 is within the leave-one-out null (0.45±0.20, p≈0.07). "
  "Per-DECISION error flagging is the real signal (AUC 0.59; survives task-difficulty control).")
prs.save(str(TALK/"berlin_deck_v8.pptx"))
print("saved berlin_deck_v8.pptx |",len(prs.slides),"slides")
