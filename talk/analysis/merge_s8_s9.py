"""Fold S8 (Can we learn a language of proofreading?) into S9 (JIGSAWS inspiration).
S9 becomes the combined slide: the proofreading-language question is the setup, JIGSAWS the answer.
S8 (the pivot) -> backup. 20 main -> 19 main, +1 backup. In place on the v14 deck.
"""
from pathlib import Path
from pptx import Presentation
DECK=Path(__file__).resolve().parent.parent/"Calibrate_the_Humans_v14.pptx"
p=Presentation(str(DECK)); S=lambda n:p.slides[n-1]

def settext(shape,text):
    tf=shape.text_frame; p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=text
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else: p0.add_run().text=text
    for para in tf.paragraphs[1:]: para._p.getparent().remove(para._p)
def setby(n,snip,new):
    for sh in S(n).shapes:
        if sh.has_text_frame and snip in sh.text_frame.text: settext(sh,new); return
    print("  !! no match S%d %r"%(n,snip[:30]))

# S9 (JIGSAWS) becomes the combined slide
setby(9,"Inspiration: JIGSAWS","A Language of Proofreading? — Inspiration: JIGSAWS (Surgery)")
setby(9,"Surgical procedures can be described as sequences of gestures and transitions",
      "Can we read proofreading skill from how the work is done — to route limited expert attention? "
      "Surgery already has a 'language': JIGSAWS describes procedures as sequences of gestures and "
      "transitions, and expert surgeons show more efficient, consistent patterns and smoother motion. "
      "Skill is legible from the grammar, not just the outcome.")
# S8 (the pivot) -> backup
setby(8,"Can we learn a language of proofreading","Backup — Can we learn a language? (pivot)")
lst=p.slides._sldIdLst; el=list(lst)[7]; lst.remove(el); lst.append(el)

p.save(str(DECK))
print("saved %s — %d slides (19 main + 5 backup)"%(DECK.name,len(p.slides._sldIdLst)))
