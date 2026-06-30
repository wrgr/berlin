"""Update the 'Risk Is Estimable, Ground-Truth-Free' slide: swap in the combined two-signal
figure (fig_task_risk.png) and rewrite the body to describe both GT-free risk signals + significance.
In place on the v14 deck."""
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
TALK=Path(__file__).resolve().parent.parent
DECK=TALK/"Calibrate_the_Humans_v14.pptx"; NEWFIG=TALK/"fig_task_risk.png"
p=Presentation(str(DECK))

def settext(shape,text):
    tf=shape.text_frame; p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=text
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else: p0.add_run().text=text
    for para in tf.paragraphs[1:]: para._p.getparent().remove(para._p)

risk=next(s for s in p.slides if any(sh.has_text_frame and "Risk Is Estimable" in sh.text_frame.text for sh in s.shapes))
# swap figure (blob replace -> preserves slot + z-order; both PNG)
img=next(sh for sh in risk.shapes if sh.shape_type==13)
img.part.related_part(img._element.blipFill.blip.get(qn("r:embed")))._blob=NEWFIG.read_bytes()
# rewrite body
body=next(sh for sh in risk.shapes if sh.has_text_frame and "Using only task structure" in sh.text_frame.text)
settext(body,
    "Two ground-truth-free ways to flag a risky decision, before any expert looks: (1) behavioral — a "
    "task slower than the annotator's OWN average is more error-prone (AUC 0.59, p<0.001); (2) structural "
    "— the task's point-category mix predicts which tasks err, AUC 0.76 on held-out cells (p<0.001; the "
    "inflated 0.92 was cell-identity leakage). Both score the risk axis of impact×risk before any expert time is spent.")
p.save(str(DECK))
print("updated risk slide; embedded fig =", len(img.image.blob), "bytes vs file", NEWFIG.stat().st_size)
