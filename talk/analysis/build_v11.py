"""Deck v10 -> v11: incorporate a review pass —
1. page numbers on every slide; 2. center the figure legends/captions; 3. NIH acknowledgment +
disclaimer on the acknowledgments slide, and JHU/APL -> APL. Scripted/reversible. BERLIN_TALK overrides."""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v10.pptx")); S=prs.slides
W=prs.slide_width; H=prs.slide_height
DISC=("The authors thank NIH’s National Institute of Neurological Disorders and Stroke for "
      "supporting this work under the “BRAIN CONNECTS: A Center for High-throughput Integrative "
      "Mouse Connectomics” award UM1NS132250. The content is solely the responsibility of the "
      "authors and does not necessarily represent the official views of NIH.")
FIGSLIDES=[13,14,15,16,17,18,23]
log=[]
# 1) page numbers (bottom-right) on every slide
for i,s in enumerate(S,1):
    tb=s.shapes.add_textbox(W-Inches(0.7),H-Inches(0.45),Inches(0.5),Inches(0.3))
    p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    r=p.add_run(); r.text=str(i); r.font.size=Pt(10); r.font.color.rgb=RGBColor(0x80,0x80,0x80)
log.append("page numbers on %d slides"%len(S))
# 2) center the figure captions (TEXT_BOX shapes, not the title placeholder, not the page-number box)
for idx in FIGSLIDES:
    for sh in S[idx-1].shapes:
        if sh.has_text_frame and sh.shape_type==17 and not sh.is_placeholder and sh.text_frame.text.strip()!=str(idx):
            for para in sh.text_frame.paragraphs: para.alignment=PP_ALIGN.CENTER
log.append("centered figure legends on slides %s"%FIGSLIDES)
# 3) acknowledgments slide (21): JHU/APL -> APL, add NIH disclaimer
s21=S[20]
for sh in s21.shapes:
    if not sh.has_text_frame: continue
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            if "JHU/APL" in run.text: run.text=run.text.replace("JHU/APL","APL"); log.append("JHU/APL -> APL")
tb=s21.shapes.add_textbox(Inches(0.6),H-Inches(1.7),Inches(11.4),Inches(1.4)); tb.text_frame.word_wrap=True
p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
r=p.add_run(); r.text=DISC; r.font.size=Pt(9); r.font.italic=True; r.font.color.rgb=RGBColor(0x55,0x55,0x55)
log.append("added NIH acknowledgment + disclaimer")
prs.save(str(TALK/"berlin_deck_v11.pptx"))
print("saved berlin_deck_v11.pptx |",len(prs.slides),"slides | edits:",log)
