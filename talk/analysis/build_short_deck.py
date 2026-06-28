#!/usr/bin/env python3
"""Build berlin_deck_short.pptx — a ~7-slide condensation of berlin_deck_v11.pptx.

Captures the main point and big idea, with one emphasis: an **agent** is agnostic to human or
machine. Experts, students, and ML are interchangeable agents in the loop; the whole talk is about
routing each decision to the agent that should make it, and pricing every agent's judgment.

Method (scripted/reversible, matching the repo's build_v* pattern): load v11, keep 7 slides, apply
targeted single-line text edits to push the agent-agnostic framing into titles/banners/close, fix
the page numbers, rewrite the speaker notes for the condensed arc, and save a NEW file. v11 is never
modified. Run with BERLIN_TALK pointing at talk/ (defaults to the repo-relative talk/ dir).

Source -> short mapping (1-based v11 slide -> short position):
  1 Title                         -> 1   (reframe: calibrate the *agents*)
  2 Problem is a decision         -> 2   (optimization, not a quality check)
  3 Route to the right agent      -> 3   (THE big idea: agent = human or machine)
  7 Allocate by impact x risk     -> 4   (where the budget goes)
  8 Can't allocate w/o measuring  -> 5   (price every agent's judgment)
 13 The evidence                  -> 6   (competence is legible, ground-truth-free)
 22 Close                         -> 7   (calibrate the agents, not just the microscope)
"""
import os
from pathlib import Path
from pptx import Presentation

HERE = Path(__file__).resolve().parent
TALK = Path(os.environ.get("BERLIN_TALK", HERE.parent))
SRC = TALK / "berlin_deck_v11.pptx"
DST = TALK / "berlin_deck_short.pptx"

# 1-based v11 slide numbers to keep, in display order.
KEEP = [1, 2, 3, 7, 8, 13, 22]

# Targeted single-line edits: orig_slide -> list of (match_substring, new_full_text).
# Matched against current shape text (first shape that contains the substring); formatting of the
# shape's first run is preserved.
EDITS = {
    1: [("Calibrate the students, not just the microscope",
         "Calibrate the agents, not just the microscope")],
    3: [("Route the work to the right agent",
         "Route every decision to the right agent — human or machine"),
        ("EASY / HARD × WHO HANDLES IT",
         "EASY / HARD × WHICH AGENT — HUMAN OR MACHINE")],
    8: [("PRICING THE BUDGET = CALIBRATING THE ANNOTATORS",
         "PRICING THE BUDGET = CALIBRATING THE AGENTS (HUMAN OR MACHINE)")],
    22: [("The missing piece is the price of human judgment.",
          "The missing piece is the price of an agent’s judgment — human or machine."),
         ("Calibrate the people, not just the microscope.",
          "Calibrate the agents — human or machine — not just the microscope.")],
}

# Condensed-arc speaker notes, keyed by orig slide number. Agent-agnostic through-line.
NOTES = {
    1: (
        "The hook and the reframe. You have a volume and a finite number of proofreaders — when do "
        "you stop, and how do you know the answer is one you can trust? I want to land one word up "
        "front: agent. In this talk an agent is whatever makes a proofreading decision — an expert, "
        "a student, or an algorithm. They're interchangeable; the whole question is which agent each "
        "decision deserves. So the title is deliberate: calibrate the agents, not just the microscope."
    ),
    2: (
        "Frame it as engineering, not quality control. We're not reconstructing everything — we're "
        "testing a specific inference, under a hard constraint: judgment is finite, expensive, and "
        "ambiguous, because agents disagree, so there's no free ground truth. That flips the question. "
        "Not 'is the connectome done?' — done is a property of the dataset. The real question is 'have "
        "we spent enough to trust the answer?' That's an optimization problem, and it's the whole talk."
    ),
    3: (
        "This is the big idea. Two axes — how hard a decision is, how much it matters. Easy and "
        "low-stakes: let a machine do it, NEURD, no review. Easy and high-stakes: automate, then "
        "verify. Hard and high-stakes: that's the precious budget. Hard and low: skip. The point I "
        "want you to leave with is the banner: experts, students, and ML are all agents in the loop. "
        "It is not human versus machine — it's matching each decision to the right agent. Everything "
        "after this is how you do that matching rigorously."
    ),
    7: (
        "Once agents are interchangeable, allocation is the game. Two axes again, but now over "
        "decisions: impact — how much a decision moves the inference — and risk — how likely it's "
        "wrong, which here means how much agents disagree. Spend your scarce budget top-right: high "
        "impact, high disagreement. High impact but already agreed — validate, don't over-edit. Low "
        "impact — make it cheap or automate. And the lever that sets the whole map: the inference "
        "under test defines what counts as impact. Change the question, the map changes."
    ),
    8: (
        "You can't allocate what you can't measure — so you have to price judgment. That means "
        "characterizing every agent: per-task, per-region competence, and how much they agree. And "
        "note this is agent-agnostic too — you price an algorithm the same way you price a person: "
        "what is it good at, where do we trust it. This is measuring the task, not judging people. It "
        "also forces an honest definition of reliable: not matching a gold standard, because there "
        "isn't one — converged agreement on the decisions the inference depends on."
    ),
    13: (
        "The evidence that this is real, not just a framework. Left: behavior alone separates expert "
        "agents from proto-expert agents — how they explore, in what order, at what tempo — at AUC up "
        "around 0.90. Right, and this is the deployable part: rank tasks by a label-free anomaly "
        "score, with no answer key, and re-reviewing the most anomalous catches errors well above "
        "chance. So competence is legible in behavior, and risky decisions can be flagged ground-"
        "truth-free. The punchline: flag risky decisions, don't rank the agents."
    ),
    22: (
        "Land it. Proofreading a connectome isn't a quality problem you finish — it's an optimization "
        "problem under a hard constraint. The missing piece has been the price of an agent's judgment "
        "— and once you can price it, per-decision and ground-truth-free, you can route every "
        "decision to the right agent, human or machine, where impact and risk are highest. That's the "
        "line to leave with: calibrate the agents, not just the microscope. Thank you."
    ),
}


def set_first_run_text(shape, new_text):
    """Replace a shape's text with new_text, preserving the first run's formatting."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = new_text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.text = new_text
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def find_shape_containing(slide, substring):
    for sh in slide.shapes:
        if sh.has_text_frame and substring in sh.text_frame.text:
            return sh
    return None


def find_pagenum_shape(slide, num):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == str(num):
            return sh
    return None


def main():
    if not SRC.exists():
        raise SystemExit(f"source deck not found: {SRC}")
    prs = Presentation(str(SRC))
    log = []

    # 1) targeted text edits (on original slide objects, before pruning)
    for orig, edits in EDITS.items():
        slide = prs.slides[orig - 1]
        for match, new in edits:
            sh = find_shape_containing(slide, match)
            if sh is None:
                print(f"!! S{orig}: no shape matched {match!r}; skipping")
                continue
            set_first_run_text(sh, new)
            log.append(f"S{orig}: edit {match[:32]!r}")

    # 2) speaker notes + 3) page numbers (new positions)
    for new_pos, orig in enumerate(KEEP, 1):
        slide = prs.slides[orig - 1]
        if orig in NOTES:
            slide.notes_slide.notes_text_frame.text = NOTES[orig]
        pg = find_pagenum_shape(slide, orig)
        if pg is not None and new_pos != orig:
            set_first_run_text(pg, str(new_pos))
            log.append(f"S{orig}: page# {orig}->{new_pos}")

    # 4) prune to the kept slides (remove others from the slide-id list)
    keep0 = {k - 1 for k in KEEP}
    id_list = prs.slides._sldIdLst
    for idx, sldId in enumerate(list(id_list)):
        if idx not in keep0:
            id_list.remove(sldId)

    prs.save(str(DST))
    print(f"saved {DST.name} | {len(prs.slides)} slides | {len(log)} edits")
    for line in log:
        print("  -", line)


if __name__ == "__main__":
    main()
