"""Deck v9 -> v10: promote the GT-free task-RISK result to a MAIN evidence slide (fig_task_risk),
inserted right after the 'per-decision, not per-person' slide as the positive capstone of the
evidence arc. Clones the kinematics figure-slide for styling. Reversible/scripted."""
import os, copy
from pathlib import Path
from pptx import Presentation
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v9.pptx"))
def shape_in(sl,sub):
    for sh in sl.shapes:
        if sh.has_text_frame and sub in sh.text_frame.text: return sh
def set_text(sh,new):
    p=sh.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text=new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else: p.add_run().text=new
    for ex in sh.text_frame.paragraphs[1:]: ex._p.getparent().remove(ex._p)
def dup(prs,i):
    src=prs.slides[i]; new=prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes): sh._element.getparent().remove(sh._element)
    for sh in src.shapes: new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new
new=dup(prs,13)                                    # clone the kinematics single-figure slide
set_text(shape_in(new,"Behavioral mechanism"),"Risk is estimable, ground-truth-free")
for sh in list(new.shapes):                        # swap fig_kinematics -> fig_task_risk (same footprint, centered)
    if sh.shape_type==13:
        l,t,w,h=sh.left,sh.top,sh.width,sh.height; sh._element.getparent().remove(sh._element)
        pic=new.shapes.add_picture(str(TALK/"fig_task_risk.png"),l,t,height=h); pic.left=l+(w-pic.width)//2; break
set_text(shape_in(new,"Experts accumulate"),
  "GT-free, from task structure: which decisions are error-prone is predictable at AUC 0.76 on held-out "
  "cells — the 'risk' axis of impact×risk, scored before any expert time is spent. (Honest grouped CV: the "
  "0.92 was cell-identity leakage; per-person competence stays ~0.55.)")
sl=prs.slides._sldIdLst; items=list(sl); sl.remove(items[-1]); sl.insert(17, items[-1])   # move to slide 18 (after 'per-decision')
prs.save(str(TALK/"berlin_deck_v10.pptx"))
print("saved berlin_deck_v10.pptx |",len(prs.slides),"slides; new evidence slide at position 18")
