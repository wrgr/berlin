"""Resample every embedded image in a deck to 300 DPI at its on-slide display size.
Downsamples oversized rasters (big file savings, still print-quality); never upsamples.
Preserves z-order (edits the image PART blob in place, not remove+add) and image format.
Usage: python resample_deck_300dpi.py [in.pptx] [out.pptx]  (defaults: berlin_deck_v12 in place)
"""
import sys, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from PIL import Image as PILImage

TALK=Path(__file__).resolve().parent.parent
SRC=Path(sys.argv[1]) if len(sys.argv)>1 else TALK/"berlin_deck_v12.pptx"
OUT=Path(sys.argv[2]) if len(sys.argv)>2 else SRC
DPI=300

prs=Presentation(str(SRC))
need={}  # image part -> (max_wpx, max_hpx) across all shapes that use it
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.shape_type!=13:  # not a picture
            continue
        rId=sh._element.blipFill.blip.get(qn("r:embed"))
        if rId is None: continue
        part=sh.part.related_part(rId)
        wpx=int(round(Emu(sh.width).inches*DPI)); hpx=int(round(Emu(sh.height).inches*DPI))
        w0,h0=need.get(part,(0,0)); need[part]=(max(w0,wpx),max(h0,hpx))

before=sum(len(p._blob) for p in need); changed=0; saved=0
for part,(wpx,hpx) in need.items():
    try:
        im=PILImage.open(io.BytesIO(part._blob)); fmt=(im.format or "PNG").upper()
    except Exception as e:
        print("  skip (unreadable): %s"%e); continue
    ratio=min(wpx/im.width, hpx/im.height)
    if ratio>=1.0:  # already <= 300 DPI at display size -> leave it (no upsampling)
        continue
    new=im.resize((max(1,int(im.width*ratio)),max(1,int(im.height*ratio))),PILImage.LANCZOS)
    buf=io.BytesIO()
    if fmt in ("JPEG","JPG"):
        if new.mode in ("RGBA","P","LA"): new=new.convert("RGB")
        new.save(buf,"JPEG",quality=88,optimize=True)
    else:
        new.save(buf,"PNG",optimize=True)
    nb=buf.getvalue()
    if len(nb)<len(part._blob):
        d=len(part._blob)-len(nb); saved+=d; changed+=1
        print("  %dx%d->%dx%d  %.2f->%.2fMB  (%s)"%(im.width,im.height,new.width,new.height,
              len(part._blob)/1e6,len(nb)/1e6,fmt))
        part._blob=nb
prs.save(str(OUT))
after=before-saved
print("\nresampled %d/%d images to <=%d DPI; image bytes %.1f -> %.1f MB (saved %.1f MB)"%(
    changed,len(need),DPI,before/1e6,after/1e6,saved/1e6))
print("saved", OUT)
