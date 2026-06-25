"""Deck v5 (your hand-edited 20-slide deck) -> v6: coherence refinements only.
Improves story flow; NO length cuts. Every edit is scripted and reversible.
See ../deck_coherence_review.md for the rationale. Run with BERLIN_TALK pointing at talk/."""
import os
from pathlib import Path
from pptx import Presentation
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v5.pptx")); S=prs.slides

def para_in(slide, substr):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            if substr in p.text: return p
    return None
def set_run0(p, new):                       # replace paragraph text, keep first run's font
    if p.runs:
        p.runs[0].text=new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else: p.add_run().text=new
def sub_run(p, old, new):                   # in-run substring replace, preserves all formatting
    for r in p.runs:
        if old in r.text: r.text=r.text.replace(old,new); return True
    return False

log=[]
# --- Slides 11 & 12 both opened with "promotion is noisy/slow/ad hoc": let 11 own the reframe,
#     make 12 ADVANCE (what behavior reveals) and bridge from 11 instead of repeating it. ---
p=para_in(S[11],"Promotion is noisy and slow — behavior shows")
if p: set_run0(p,"What behavior already reveals"); log.append("S12 retitle (was a repeat of S11)")
p=para_in(S[11],"Today, promoting a proofreader is noisy, hard, and slow")
if p: set_run0(p,"Skill lives in the how — so we can read it as the work happens, instead of waiting on outcomes."); log.append("S12 opening: bridge from S11, not a repeat")
p=para_in(S[11],"But competence is legible in behavior")
if p and sub_run(p,"But competence is legible","Competence is legible"): log.append("S12 drop leading 'But'")
# --- Two impact-grids (3 & 7) read as near-duplicates: distinguish slide 7 + bridge in the notes ---
p=para_in(S[6],"THE STRUCTURE OVER DECISIONS")
if p: set_run0(p,"THE STRUCTURE OVER DECISIONS — WHERE THE BUDGET GOES"); log.append("S7 subtitle distinguishes it from S3 ('who handles it')")
if S[6].has_notes_slide:
    ntf=S[6].notes_slide.notes_text_frame
    ntf.text=ntf.text.rstrip()+" (Bridge from the earlier grid: slide 3 routed each decision to the right agent — human or machine — by difficulty; here we spend the human budget itself, crossing impact with risk / disagreement.)"
    log.append("S7 notes: bridge to S3")
# --- Slide 17 caption read backwards vs its own figure axis (distance-to-GT, lower = better) ---
p=para_in(S[16],"Promoted ≈ expert < unpromoted")
if p: set_run0(p,"Promoted ≈ expert < unpromoted — split distance-to-GT (lower = better)."); log.append("S17 caption matches the figure axis")
prs.save(str(TALK/"berlin_deck_v6.pptx"))
print("saved berlin_deck_v6.pptx |",len(prs.slides),"slides | edits:",log)
