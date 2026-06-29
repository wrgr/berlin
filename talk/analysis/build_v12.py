"""Deck v11 -> v12: thread the point-agreement result into the evidence arc.

1. Footnote on slide 8 ("converged agreement", measured) and slide 17 ("per-decision, not
   per-person") with the new number: calibrated annotators match the expert grader (Pat) on ~99%
   of individual point decisions, while the unpromoted group is bimodal with a low tail.
2. A BACKUP slide (end of deck) showing fig_point_agreement.png — cloned from the kinematics
   figure-slide for styling, image swapped, title + caption set, page number fixed.
Companion figure: ../fig_point_agreement.png (make_point_agreement_figure.py). Scripted/reversible;
the footnotes only ADD textboxes and the slide is appended, so nothing existing is disturbed.
Run with BERLIN_TALK pointing at talk/.
"""
import os, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = Path(__file__).resolve().parent
TALK = Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs = Presentation(str(TALK / "berlin_deck_v11.pptx")); S = prs.slides

def add_footnote(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.62), Inches(11.4), Inches(0.36))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    f = r.font; f.size = Pt(11); f.italic = True; f.color.rgb = RGBColor(0x2A, 0x4D, 0x7A)
    return tb

def shape_in(sl, sub):
    for sh in sl.shapes:
        if sh.has_text_frame and sub in sh.text_frame.text: return sh

def set_text(sh, new):
    p = sh.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else: p.add_run().text = new
    for ex in sh.text_frame.paragraphs[1:]: ex._p.getparent().remove(ex._p)

def dup(prs, i):
    src = prs.slides[i]; new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes): sh._element.getparent().remove(sh._element)
    for sh in src.shapes: new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new

log = []
add_footnote(S[7], "Converged agreement, measured: calibrated annotators match the expert grader on "
                   "~99% of individual point decisions (promoted median 100%, expert 98%).  → fig_point_agreement")
log.append("S8 footnote: ~99% point-agreement (converged agreement)")
add_footnote(S[16], "Per-decision convergence: promoted ≈ expert — both ~99% point-agreement with the grader; "
                    "the unpromoted group is bimodal with a low tail.  → fig_point_agreement")
log.append("S17 footnote: promoted ≈ expert ~99%, unpromoted bimodal")

# Backup slide: fig_point_agreement (clone the kinematics single-figure slide for styling)
new = dup(prs, 13)
set_text(shape_in(new, "Behavioral mechanism"), "Backup — calibration converges, per decision (point agreement)")
for sh in list(new.shapes):                                  # swap kinematics image -> fig_point_agreement (wide; fit width, centered)
    if sh.shape_type == 13:
        t = sh.top; sh._element.getparent().remove(sh._element)
        pic = new.shapes.add_picture(str(TALK / "fig_point_agreement.png"), Inches(0.9), t, width=Inches(11.5))
        pic.left = (prs.slide_width - pic.width) // 2; break
set_text(shape_in(new, "Experts accumulate"),
  "Per-point label agreement with the expert grader (Pat) on the shared benchmark cells: promoted ≈ "
  "expert (medians 100% / 98%), unpromoted bimodal with a low tail. Representative expert 99.3% "
  "(141/142). Decisions converge even where behavioral style differs — style ≠ proficiency.")
pg = next((sh for sh in new.shapes if sh.has_text_frame and sh.text_frame.text.strip() == "14"), None)
if pg is not None: set_text(pg, str(len(prs.slides)))         # fix cloned page number -> last slide
new.notes_slide.notes_text_frame.text = (
    "Backup — only if asked 'how close to ground truth are your annotators?'. These are forced-choice "
    "point-classification tasks: the same points are pre-placed, everyone labels them, so I can measure "
    "per-point agreement with the expert grader directly. Promoted students match the grader on a median "
    "100% of points, experts 98%; a representative expert is 141 of 142. The unpromoted group is bimodal "
    "with a low tail — which is exactly why you flag decisions, not rank people. And the punchline: this "
    "decision-agreement is decoupled from behavioral style — the same expert agrees ~99% on labels but "
    "works in a totally different behavioral dialect. Style is manner; the points are correctness.")
log.append(f"backup slide appended at position {len(prs.slides)} (fig_point_agreement)")

prs.save(str(TALK / "berlin_deck_v12.pptx"))
print("saved berlin_deck_v12.pptx |", len(prs.slides), "slides | edits:", log)
