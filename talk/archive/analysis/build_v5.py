"""Deck v4 -> v5 edits (scripted so every change is reversible; see ../deck_v5_changes.md).
- Slide 8: add INFERENCE as the mediator of impact (title, y-axis, banner, takeaway).
- Slide 15: bind 'language of proofreading' to the evidence title.
- Slide 7: header order matches body (Correction then Validation).
Run from anywhere; paths are repo-relative (override base/out with BERLIN_TALK)."""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))
prs=Presentation(str(TALK/"berlin_deck_v4.pptx")); S=prs.slides

def set_text(shape,new):
    """Replace a single-paragraph shape's text, preserving the first run's font."""
    tf=shape.text_frame; p=tf.paragraphs[0]
    if p.runs:
        p.runs[0].text=new
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else:
        p.add_run().text=new
    for ex in tf.paragraphs[1:]: ex._p.getparent().remove(ex._p)

log=[]
# ---- Slide 8: INFERENCE sets impact ----
for sh in S[7].shapes:
    if not sh.has_text_frame: continue
    t=sh.text_frame.text.strip()
    if t=="Allocate by impact × risk":
        set_text(sh,"Allocate by impact × risk — set by the inference"); log.append("S8 title")
    elif t=="impact ↑":
        set_text(sh,"impact on the inference ↑"); log.append("S8 y-axis")
    elif t.startswith("Route each decision"):
        set_text(sh,"Route each decision — and the right annotator — to where impact and risk are both high. Impact is defined by the inference you're testing."); log.append("S8 takeaway")
tb=S[7].shapes.add_textbox(Inches(0.8),Inches(1.63),Inches(11.6),Inches(0.34))
r=tb.text_frame.paragraphs[0].add_run()
r.text="▲  The inference under test sets the impact axis — the hypothesis mediates what's worth the budget."
r.font.size=Pt(13); r.font.italic=True; r.font.color.rgb=RGBColor(0x2C,0x7F,0x8C); log.append("S8 banner")
# ---- Slide 15: language of proofreading ----
for sh in S[14].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip()=="The evidence (handles suppressed)":
        set_text(sh,"The evidence — the language of proofreading (handles suppressed)"); log.append("S15 title")
# ---- Slide 7: header order ----
for sh in S[6].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip()=="VALIDATION VS. CORRECTION":
        set_text(sh,"CORRECTION VS. VALIDATION"); log.append("S7 header")
prs.save(str(TALK/"berlin_deck_v5.pptx"))
print("saved berlin_deck_v5.pptx |",len(prs.slides),"slides | edits:",log)
