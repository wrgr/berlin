"""Expand berlin_deck_v2.pptx -> berlin_deck_v3.pptx: add cohort, learning-engineering,
evidence, outreach, acknowledgments slides; reorder into the narrative."""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
HERE=Path(__file__).resolve().parent
TALK=Path(os.environ.get("BERLIN_TALK", HERE.parent))                       # repo talk/
SRC=os.environ.get("BERLIN_DECK_SRC", str(TALK/"berlin_deck_v2.pptx"))      # base deck
DST=str(TALK/"berlin_deck_v3.pptx")
prs=Presentation(SRC)
# pick a content layout: reuse an existing content slide's layout
layout=prs.slides[2].slide_layout
SW=prs.slide_width; SH=prs.slide_height
AST=str(TALK)+"/assets/"
TALK_FIG=str(TALK)+"/"
def add(title,kicker,bullets,bottom,notes,caveat=None,img=None,imgcap=None):
    s=prs.slides.add_slide(layout)
    # title
    tph=None; body=None
    for ph in s.placeholders:
        if ph.placeholder_format.idx==0: tph=ph
        elif body is None: body=ph
    if tph is not None: tph.text=title
    else:
        tb=s.shapes.add_textbox(Inches(0.6),Inches(0.3),Inches(9),Inches(0.9)); tb.text_frame.text=title; tph=tb
    # body
    if body is None:
        body=s.shapes.add_textbox(Inches(0.7),Inches(1.5),Inches(8.7),Inches(5.2))
    if img:
        try: body.left=Inches(0.5); body.width=int(SW*0.55)
        except Exception: pass
    tf=body.text_frame; tf.word_wrap=True; tf.clear()
    p=tf.paragraphs[0]; p.text=kicker
    for r in p.runs: r.font.bold=True; r.font.size=Pt(13); r.font.color.rgb=RGBColor(0x2E,0x74,0xB5)
    for b in bullets:
        p=tf.add_paragraph(); p.text=b; p.level=1
        for r in p.runs: r.font.size=Pt(14 if img else 15)
    if bottom:
        p=tf.add_paragraph(); p.text=bottom
        for r in p.runs: r.font.bold=True; r.font.size=Pt(14 if img else 15); r.font.color.rgb=RGBColor(0x1F,0x49,0x77)
    if caveat:
        p=tf.add_paragraph(); p.text=caveat
        for r in p.runs: r.font.italic=True; r.font.size=Pt(11); r.font.color.rgb=RGBColor(0x90,0x40,0x40)
    if img:
        pl=int(SW*0.58); pt=int(SH*0.20); pic=s.shapes.add_picture(img,pl,pt)
        sc=min(int(SW*0.40)/pic.width,int(SH*0.66)/pic.height); pic.width=int(pic.width*sc); pic.height=int(pic.height*sc)
        if imgcap:
            cap=s.shapes.add_textbox(pl,pt+pic.height,pic.width,Inches(0.5)); cap.text_frame.word_wrap=True; cap.text_frame.text=imgcap
            for r in cap.text_frame.paragraphs[0].runs: r.font.size=Pt(9); r.font.italic=True
    # notes
    s.notes_slide.notes_text_frame.text=notes
    return s
def add_figslide(title,imgs,caps,notes):
    s=prs.slides.add_slide(layout)
    for ph in s.placeholders:
        if ph.placeholder_format.idx==0: ph.text=title
        else:
            try: ph.text_frame.clear()
            except Exception: pass
    n=len(imgs); margin=int(SW*0.04); gap=int(SW*0.03)
    boxw=int((SW-2*margin-(n-1)*gap)/n); top=int(SH*0.24); boxh=int(SH*0.56); x=margin
    for im,cp in zip(imgs,caps):
        pic=s.shapes.add_picture(im,x,top); sc=min(boxw/pic.width,boxh/pic.height)
        pic.width=int(pic.width*sc); pic.height=int(pic.height*sc); pic.left=x+int((boxw-pic.width)/2)
        cap=s.shapes.add_textbox(x,top+pic.height,boxw,Inches(0.5)); cap.text_frame.word_wrap=True; cap.text_frame.text=cp
        for r in cap.text_frame.paragraphs[0].runs: r.font.size=Pt(10); r.font.italic=True
        x+=boxw+gap
    s.notes_slide.notes_text_frame.text=notes
    return s
add("Who we calibrated: the proofreading workforce","THE PEOPLE BEHIND THE BUDGET",
    ["8 part-time expert proofreaders — prior EM / neuroanatomy training at Janelia (Nov 2021 – Sep 2022).",
     "36 novice proofreaders — Johns Hopkins undergraduates: 26 founders (3 weeks training, then full-time in January, then part-time through August); 10 more joined in June 2022.",
     "One curriculum: neuroanatomy, EM interpretation, segmentation overlays, the NeuVue interface, and the decision logic for each task type.",
     "Routed and recorded on NeuVue (Xenes et al.)."],
    "“Calibrate the annotators” isn’t a metaphor — it’s 44 people and a training pipeline.",
    "Grounds the abstract budget in the real workforce. Next: the calibration mechanism, and what behavior reveals about it.")
add("Learning engineering: agreement-gated promotion","CALIBRATION WAS DESIGNED, NOT ASSUMED",
    ["Students whose decisions agreed with experts were promoted to a proto-expert tier — with write permission for expert-level tasks.",
     "A learning-engineering loop: train → practice → measure agreement → promote → widen scope.",
     "Proto-expert cohort: dylan, vivia, taylor, clara, rachel, shruthi, sarah, lydia.",
     "Reused in CONNECTS-Proof to validate NEURD edits on H01.",
     "Same thesis as APL’s MERIT / AI Pathfinding: assessment as a learning diagnostic; deliverables that outlive the cohort."],
    "Promotion-after-agreement is selection on outcome — can we predict it from early behavior instead?",
    "Bridges to the original staged-promotion slide and to the predictive test. The promoted list is authoritative (project lead).",
    img=AST+"incubator_lab.jpg",imgcap="Mentored research in practice — APL Research Incubator")
add("Preliminary evidence: competence is legible in behavior","THE LANGUAGE OF PROOFREADING",
    ["Dense per-event telemetry: every navigate / edit / annotate, with 3D camera motion.",
     "Behavior separates experts from proto-experts: naive 0.75 → designed 0.95 → learned motif dictionary 0.90.",
     "Top signals are 3D-exploration kinematics — experts inspect from ~2× more viewpoints.",
     "Yet achieved accuracy is ceiling-clustered across two task types — a calibrated workforce converges on outcomes.",
     "Prospective & GT-free: a task slow FOR THAT PERSON is significantly more error-prone (AUC 0.59, p<0.001) — ‘subconscious uncertainty.’",
     "The language of proofreading ↔ the language of surgery (JIGSAWS; tacit knowledge): skill lives in HOW, not just WHAT."],
    "Expertise is a behavioral style; per-task uncertainty stays legible even after competence converges.",
    "Style separation is retrospective; the per-task GT-free signal is small (AUC ~0.59) and the prospective test is this summer. Honest negative: annotator-level simple-behavior separability FAILED (LOO AUC 0.14).")
add("Outreach: students as a method, not just a mission","HUMANS + OUTCOMES   (LENS)",
    ["Calibrated student judgment IS the method that built the connectome — not a side benefit.",
     "MICrONS relied on this workforce throughout — their proofreading underpins the MICrONS connectome (Nature, 2025).",
     "Students we hired became co-authors and contributing members — Daniel Xenes (now leading NeuVue) is the exemplar.",
     "Several novices were top-throughput proofreaders credited in MICrONS; eight were promoted to expert tasks.",
     "APL NeuroTrailblazers / Research Incubator lineage; MERIT learning engineering (assess → mentor → promote).",
     "Pricing judgment is what makes “students as a method” rigorous and scalable."],
    "The measurement that prices the budget is the same measurement that develops the workforce.",
    "Dual yield is structurally identical — talent development and mission acceleration are the same act. LENS = humans + outcomes; students as a METHOD (Gray-Roncal, Research Incubator).",
    img=AST+"incubator_cohort.jpg",imgcap="APL Research Incubator cohort — students as a method")
add("Acknowledgments & contributions","A TEAM EFFORT",
    ["NeuVue platform & task routing — Daniel Xenes and the NeuVue team (the substrate for all of this).",
     "Graders / experts — Pat Rivlin, Lindsey Kitchell; the 8-person expert cohort.",
     "Proofreading team — the 36 JHU novices; top acknowledged contributors incl. N. Smith (Natalie, 24,101 edits), D. Panchal, M. Cook, C. Ordish, C. Knecht, E. Phillips.",
     "Datasets & methods — MICrONS Consortium (Nature, 2025); Matelsky & Gray-Roncal (consequence-of-edits / motif census); Sanchez et al., Connectomics Annotation Metadata Standardization (Front. Neuroinform. 2022); H01 / CONNECTS-Proof.",
     "Workforce model — APL NeuroTrailblazers / Research Incubator + MERIT (Cervantes, Floryanzia, Johnson, Gray-Roncal et al.).",
     "Institutions — JHU APL · Johns Hopkins University · Janelia Research Campus · Allen Institute."],
    "",
    "Daniel Xenes especially, and the full team. Sanchez ref = Front. Neuroinform. 16:828458 (2022).")
add_figslide("Backup — the evidence (handles suppressed)",
    [TALK_FIG+"fig_tier_auc.png",TALK_FIG+"fig_prospective_flagging.png"],
    ["Behavior separates experts from proto-experts (naive→designed→learned).",
     "Ground-truth-free: flag the most behaviorally-anomalous tasks → catch errors above chance."],
    "Two data results: expertise is legible in behavioral style; individual risky decisions can be flagged without ground truth (the prospective hook).")
add_figslide("Backup — behavioral mechanism (3-D exploration & tempo)",
    [TALK_FIG+"fig_kinematics.png"],
    ["Experts accumulate ~2× more camera rotation, inspect from more viewpoints, do more, and move faster."],
    "Mechanism behind the AUC: expertise is a 3-D exploration style.")
add_figslide("Backup — the language of proofreading (action mix & grammar)",
    [TALK_FIG+"fig_action_grammar.png"],
    ["Action mix and navigate↔segment transition grammar differ by cohort."],
    "The behavioral 'language': which actions, in what order.")
add_figslide("Backup — what separates experts from proto-experts",
    [TALK_FIG+"fig_rf_importance_new.png",TALK_FIG+"fig_feature_pca.png",TALK_FIG+"fig_motif_usage.png"],
    ["RandomForest importance (designed features).","Designed feature space (PCA).","Learned-motif usage (dialect)."],
    "Importance, low-dim structure, and learned-motif dialect all point to exploration kinematics.")
add_figslide("Backup — calibration converges; uncertainty stays legible",
    [TALK_FIG+"fig_accuracy_threegroup.png",TALK_FIG+"fig_uncertainty_calibration.png"],
    ["Promoted ≈ expert < unpromoted accuracy.","Slower-for-this-person tasks fail more (GT-free)."],
    "Calibration compresses outcome variance; the surviving signal is per-decision uncertainty.")
# reorder: A,B after slide8(idx7); orig8 stays; C,D,E before closing(orig9); figslides last (backup)
sldIdLst=prs.slides._sldIdLst; els=list(sldIdLst)
order=[0,1,2,3,4,5,6,7,10,11,8,12,13,14,9,15,16,17,18,19]
for e in els: sldIdLst.remove(e)
for i in order: sldIdLst.append(els[i])
prs.save(DST)
print("saved",DST)
# verify
p2=Presentation(DST)
print("verify N slides:",len(p2.slides))
for i,s in enumerate(p2.slides):
    t=""
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t=sh.text_frame.text.strip().split("\n")[0]; break
    print(" %2d: %s"%(i+1,t[:60]))
